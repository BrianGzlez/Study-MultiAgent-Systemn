from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import io


def extract_text(file_bytes: bytes, file_type: str) -> str:
    """Extract text from uploaded files based on file type."""
    file_type = file_type.lower()

    if file_type == "pdf":
        return _extract_from_pdf(file_bytes)
    elif file_type == "docx":
        return _extract_from_docx(file_bytes)
    elif file_type == "pptx":
        return _extract_from_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def _extract_from_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text)
    return "\n\n".join(texts)


def _extract_from_docx(file_bytes: bytes) -> str:
    doc = DocxDocument(io.BytesIO(file_bytes))
    texts = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            texts.append(paragraph.text)
    return "\n\n".join(texts)


def _extract_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    texts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            texts.append(" ".join(slide_text))
    return "\n\n".join(texts)


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> list[str]:
    """Split text into overlapping chunks."""
    sentences = text.replace("\n", " ").split(". ")
    chunks = []
    current_chunk = ""

    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue

        if len(current_chunk) + len(sentence) > chunk_size and current_chunk:
            chunks.append(current_chunk.strip())
            # Keep some overlap
            words = current_chunk.split()
            overlap_words = words[-overlap // 5:] if len(words) > overlap // 5 else words
            current_chunk = " ".join(overlap_words) + ". " + sentence
        else:
            current_chunk += (". " if current_chunk else "") + sentence

    if current_chunk.strip():
        chunks.append(current_chunk.strip())

    return [c for c in chunks if len(c) > 50]
